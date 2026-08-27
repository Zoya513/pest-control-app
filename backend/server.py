"""PestOps Pro - Field Operations Management Platform for Pest Control.
Backend: FastAPI + MongoDB.
"""
from dotenv import load_dotenv
from pathlib import Path
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / ".env")

import os
import io
import uuid
import math
import logging
from datetime import datetime, timezone, timedelta, date
from typing import Optional, List, Dict, Any

from fastapi import FastAPI, APIRouter, HTTPException, Depends, Request, Response, UploadFile, File, Query, Header, Body
from fastapi.responses import StreamingResponse
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel, Field, EmailStr

from auth import (
    hash_password, verify_password, create_access_token, get_current_user,
    require_permission, has_permission, default_admin_permissions,
    default_technician_permissions, default_client_permissions,
    default_developer_permissions, MODULES, ACTIONS
)
from storage import init_storage, put_object, get_object, APP_NAME
from reports import (
    generate_service_report_pdf, generate_attendance_excel,
    generate_simple_pdf, generate_simple_excel
)
from emailer import send_email_with_attachments, render_report_email, EMAIL_FROM_NAME
from smtp_sender import send_via_smtp, render_template, body_to_html, DEFAULT_SR_SUBJECT, DEFAULT_SR_BODY, DEFAULT_MR_SUBJECT, DEFAULT_MR_BODY
from wa_sender import send_whatsapp, DEFAULT_SR_WA, DEFAULT_MR_WA
import httpx
import zipfile
import csv
import io as _io

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("pestops")

mongo_url = os.environ["MONGO_URL"]
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ["DB_NAME"]]

app = FastAPI(title="PestOps Pro")
api = APIRouter(prefix="/api")


# ---------- utilities ----------
def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def uid() -> str:
    return str(uuid.uuid4())


def clean(doc: dict) -> dict:
    if not doc:
        return doc
    doc.pop("_id", None)
    doc.pop("password_hash", None)
    return doc


def haversine_m(lat1, lon1, lat2, lon2) -> float:
    R = 6371000.0
    p1, p2 = math.radians(lat1), math.radians(lat2)
    dp = math.radians(lat2 - lat1)
    dl = math.radians(lon2 - lon1)
    a = math.sin(dp / 2) ** 2 + math.cos(p1) * math.cos(p2) * math.sin(dl / 2) ** 2
    return 2 * R * math.asin(math.sqrt(a))


async def audit(user: dict, action: str, module: str, record_id: str, old=None, new=None, ip: str = ""):
    await db.audit_logs.insert_one({
        "id": uid(),
        "user_id": user["id"],
        "user_name": user.get("full_name") or user["email"],
        "action": action,
        "module": module,
        "record_id": record_id,
        "old_value": old,
        "new_value": new,
        "ip": ip,
        "timestamp": now_iso(),
    })


# ---------- Models (pydantic) ----------
class LoginBody(BaseModel):
    email: EmailStr
    password: str


class UserCreate(BaseModel):
    email: EmailStr
    password: str
    full_name: str
    role: str = "technician"  # admin | technician
    position: Optional[str] = ""
    phone: Optional[str] = ""
    address: Optional[str] = ""
    id_number: Optional[str] = ""
    leave_quota: int = 12


class UserUpdate(BaseModel):
    full_name: Optional[str] = None
    position: Optional[str] = None
    phone: Optional[str] = None
    address: Optional[str] = None
    id_number: Optional[str] = None
    leave_quota: Optional[int] = None
    status: Optional[str] = None
    role: Optional[str] = None
    permissions: Optional[dict] = None
    profile_photo: Optional[str] = None


class CustomerCreate(BaseModel):
    company_name: str
    contact_person: str = ""
    phone: str = ""
    email: str = ""
    project_name: str = ""
    address: str = ""
    location_text: str = ""
    latitude: Optional[float] = None
    longitude: Optional[float] = None
    category: str = "Regular"
    contract_start: Optional[str] = None
    contract_end: Optional[str] = None
    photo: Optional[str] = None
    # Optional client login credentials (when Admin wants to create login-enabled client account)
    client_email: Optional[str] = None
    client_password: Optional[str] = None


class TaskCreate(BaseModel):
    customer_id: str
    technician_id: str
    scheduled_date: str  # YYYY-MM-DD
    scheduled_time: str = "09:00"
    deadline: Optional[str] = None  # ISO
    work_target: str
    work_description: str = ""


class GPSPing(BaseModel):
    latitude: float
    longitude: float
    accuracy: Optional[float] = None
    speed: Optional[float] = None
    heading: Optional[float] = None
    task_id: Optional[str] = None
    device_id: Optional[str] = ""


class PestFinding(BaseModel):
    code: str  # F/M/C/R/A/O
    description: str = ""
    quantity: int = 0


class SRPhoto(BaseModel):
    path: str
    caption: str = ""


class SRTreatment(BaseModel):
    name: str
    area_description: str = ""


class ServiceReportCreate(BaseModel):
    task_id: str
    pest_description: str = ""
    scope_of_area: str = ""
    service_area: str = ""
    recommendation: str = ""
    pest_findings: List[PestFinding] = []
    service_treatments: List[SRTreatment] = []
    technician_signature: Optional[str] = None  # storage path
    client_signature: Optional[str] = None
    photos: List[SRPhoto] = []


class LeaveCreate(BaseModel):
    leave_type: str = "Cuti"  # Cuti | Izin
    start_date: str
    end_date: str
    return_date: Optional[str] = None
    start_time: Optional[str] = None
    reason: str
    photo: Optional[str] = None


class AttendanceCheckIn(BaseModel):
    latitude: float
    longitude: float
    accuracy: Optional[float] = None
    photo: str  # storage path
    task_id: Optional[str] = None


class AttendanceCheckOut(BaseModel):
    latitude: float
    longitude: float
    accuracy: Optional[float] = None
    photo: str
    attendance_id: str


class HeartbeatBody(BaseModel):
    gps_status: str = "active"  # active|searching|limited|denied
    latitude: Optional[float] = None
    longitude: Optional[float] = None


class ScheduleMassCreate(BaseModel):
    customer_id: str
    technician_id: str
    start_date: str  # YYYY-MM-DD
    end_date: str
    start_time: str = "08:00"
    end_time: str = "17:00"
    weekdays: List[int] = [0, 1, 2, 3, 4, 5]  # 0=Mon..6=Sun
    notes: str = ""


class ScheduleUpdate(BaseModel):
    technician_id: Optional[str] = None
    start_time: Optional[str] = None
    end_time: Optional[str] = None
    status: Optional[str] = None
    notes: Optional[str] = None


class SendReportEmail(BaseModel):
    subject: str
    message: str = ""
    override_recipient: Optional[EmailStr] = None


# ---------- geocoding helpers (Nominatim proxy) ----------
async def reverse_geocode(lat: float, lon: float) -> str:
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            r = await client.get(
                "https://nominatim.openstreetmap.org/reverse",
                params={"lat": lat, "lon": lon, "format": "json", "zoom": 18, "addressdetails": 1},
                headers={"User-Agent": "PestOpsPro/1.0"},
            )
            r.raise_for_status()
            return r.json().get("display_name", "")
    except Exception:
        return ""


# ---------- Startup / Seed ----------
@app.on_event("startup")
async def startup():
    # Indexes
    await db.users.create_index("email", unique=True)
    await db.users.create_index("id", unique=True)
    await db.customers.create_index("id", unique=True)
    await db.tasks.create_index("id", unique=True)
    await db.tasks.create_index("technician_id")
    await db.tasks.create_index("customer_id")
    await db.gps_tracking.create_index([("user_id", 1), ("timestamp", -1)])
    await db.gps_tracking.create_index("task_id")
    await db.attendance.create_index([("user_id", 1), ("date", -1)])
    await db.service_reports.create_index("task_id")
    await db.leave_requests.create_index("user_id")
    await db.audit_logs.create_index([("timestamp", -1)])
    await db.notifications.create_index([("user_id", 1), ("created_at", -1)])

    # Storage init
    init_storage()

    # Seed admin
    admin_email = os.environ.get("ADMIN_EMAIL", "admin@pestops.com")
    admin_password = os.environ.get("ADMIN_PASSWORD", "Admin@123")
    existing = await db.users.find_one({"email": admin_email})
    if not existing:
        await db.users.insert_one({
            "id": uid(),
            "email": admin_email,
            "password_hash": hash_password(admin_password),
            "full_name": "System Administrator",
            "role": "admin",
            "position": "Administrator",
            "phone": "",
            "address": "",
            "id_number": "",
            "leave_quota": 12,
            "leave_used": 0,
            "status": "active",
            "profile_photo": None,
            "permissions": default_admin_permissions(),
            "created_at": now_iso(),
            "updated_at": now_iso(),
        })
        logger.info(f"Seeded admin {admin_email}")
    else:
        # Sync password if changed
        if not verify_password(admin_password, existing.get("password_hash", "")):
            await db.users.update_one({"email": admin_email}, {"$set": {"password_hash": hash_password(admin_password)}})

    # Seed demo technician (re-sync password & default permissions on every startup)
    tech_email = "technician@pestops.com"
    tech_existing = await db.users.find_one({"email": tech_email})
    if not tech_existing:
        await db.users.insert_one({
            "id": uid(),
            "email": tech_email,
            "password_hash": hash_password("Tech@123"),
            "full_name": "Wawan Gunawan",
            "role": "technician",
            "position": "Field Technician",
            "phone": "+62 812-3456-7890",
            "address": "Depok, Jawa Barat",
            "id_number": "3276012345670001",
            "leave_quota": 12,
            "leave_used": 0,
            "status": "active",
            "profile_photo": None,
            "permissions": default_technician_permissions(),
            "created_at": now_iso(),
            "updated_at": now_iso(),
        })
        logger.info("Seeded demo technician")
    else:
        # Re-sync default permissions for seeded technician to prevent drift
        await db.users.update_one(
            {"email": tech_email},
            {"$set": {"permissions": default_technician_permissions(), "role": "technician", "status": "active"}}
        )

    # Seed developer
    dev_email = "developer@pestops.com"
    if not await db.users.find_one({"email": dev_email}):
        await db.users.insert_one({
            "id": uid(),
            "email": dev_email,
            "password_hash": hash_password("Dev@123"),
            "full_name": "Developer Account",
            "role": "developer",
            "position": "System Developer",
            "phone": "", "address": "", "id_number": "",
            "leave_quota": 0, "leave_used": 0,
            "status": "active", "profile_photo": None,
            "permissions": default_developer_permissions(),
            "created_at": now_iso(), "updated_at": now_iso(),
        })
        logger.info("Seeded developer")
    else:
        await db.users.update_one({"email": dev_email},
                                  {"$set": {"permissions": default_developer_permissions(), "role": "developer", "status": "active"}})

    # Seed demo customer + client login
    demo_customer_id = None
    demo_cust = await db.customers.find_one({"company_name": "PT. John Robert Powers"})
    if not demo_cust:
        demo_customer_id = uid()
        await db.customers.insert_one({
            "id": demo_customer_id,
            "company_name": "PT. John Robert Powers",
            "project_name": "JRP Kelapa Gading Office",
            "contact_person": "ahmad",
            "phone": "+62-21-1234567",
            "email": "client@pestops.com",
            "address": "Jl. Boulevard Raya No.1 LA 3, Kelapa Gading, Jakarta Utara",
            "location_text": "Kelapa Gading",
            "latitude": -6.1568, "longitude": 106.9051,
            "category": "Corporate",
            "contract_start": "2026-01-01", "contract_end": "2026-12-31",
            "status": "active",
            "registration_date": now_iso(),
            "created_by": "system",
            "created_at": now_iso(), "updated_at": now_iso(),
        })
    else:
        demo_customer_id = demo_cust["id"]

    # Seed client login
    client_email = "client@pestops.com"
    if not await db.users.find_one({"email": client_email}):
        await db.users.insert_one({
            "id": uid(),
            "email": client_email,
            "password_hash": hash_password("Client@123"),
            "full_name": "PT. John Robert Powers",
            "role": "client",
            "position": "Client",
            "customer_id": demo_customer_id,
            "phone": "", "address": "",
            "id_number": "",
            "leave_quota": 0, "leave_used": 0,
            "status": "active", "profile_photo": None,
            "permissions": default_client_permissions(),
            "created_at": now_iso(), "updated_at": now_iso(),
        })
        logger.info("Seeded client account")
    else:
        await db.users.update_one({"email": client_email},
                                  {"$set": {"permissions": default_client_permissions(),
                                            "role": "client", "customer_id": demo_customer_id, "status": "active"}})

    # Update credentials file
    try:
        creds_path = Path("/app/memory/test_credentials.md")
        creds_path.parent.mkdir(parents=True, exist_ok=True)
        creds_path.write_text(
            "# Test Credentials — PestOps Pro\n\n"
            "## Admin\n"
            f"- Email: `{admin_email}`\n"
            f"- Password: `{admin_password}`\n"
            "- Role: admin (full permissions)\n\n"
            "## Technician\n"
            "- Email: `technician@pestops.com`\n"
            "- Password: `Tech@123`\n"
            "- Role: technician (default field permissions)\n\n"
            "## Client\n"
            "- Email: `client@pestops.com`\n"
            "- Password: `Client@123`\n"
            "- Role: client (scoped to their own customer_id)\n\n"
            "## Developer\n"
            "- Email: `developer@pestops.com`\n"
            "- Password: `Dev@123`\n"
            "- Role: developer (branding/settings only)\n\n"
            "## Auth endpoints\n"
            "- POST `/api/auth/login`\n"
            "- POST `/api/auth/logout`\n"
            "- GET  `/api/auth/me`\n"
        )
    except Exception as e:
        logger.warning(f"credentials file write: {e}")


# ================= AUTH =================
@api.post("/auth/login")
async def login(body: LoginBody, response: Response, request: Request):
    user = await db.users.find_one({"email": body.email.lower()})
    if not user or not verify_password(body.password, user.get("password_hash", "")):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    if user.get("status") != "active":
        raise HTTPException(status_code=403, detail="Account disabled")
    token = create_access_token(user["id"], user["email"])
    response.set_cookie("access_token", token, httponly=True, secure=True, samesite="none", max_age=60 * 60 * 12, path="/")
    return {"user": clean(user), "token": token}


@api.post("/auth/logout")
async def logout(response: Response):
    response.delete_cookie("access_token", path="/")
    return {"ok": True}


@api.get("/auth/me")
async def me(user: dict = Depends(get_current_user)):
    return clean(user)


# ================= USERS =================
@api.get("/users")
async def list_users(q: Optional[str] = None, user: dict = Depends(get_current_user)):
    require_permission(user, "members", "view")
    query = {}
    if q:
        query["$or"] = [
            {"email": {"$regex": q, "$options": "i"}},
            {"full_name": {"$regex": q, "$options": "i"}},
        ]
    users = await db.users.find(query, {"password_hash": 0, "_id": 0}).to_list(500)
    return users


@api.post("/users")
async def create_user(body: UserCreate, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "members", "create")
    if await db.users.find_one({"email": body.email.lower()}):
        raise HTTPException(400, "Email already exists")
    perms = default_admin_permissions() if body.role == "admin" else default_technician_permissions()
    doc = {
        "id": uid(),
        "email": body.email.lower(),
        "password_hash": hash_password(body.password),
        "full_name": body.full_name,
        "role": body.role,
        "position": body.position or "",
        "phone": body.phone or "",
        "address": body.address or "",
        "id_number": body.id_number or "",
        "leave_quota": body.leave_quota,
        "leave_used": 0,
        "status": "active",
        "profile_photo": None,
        "permissions": perms,
        "created_at": now_iso(),
        "updated_at": now_iso(),
    }
    await db.users.insert_one(doc)
    await audit(user, "CREATE", "members", doc["id"], None, {"email": doc["email"]}, request.client.host if request.client else "")
    return clean(doc)


@api.get("/users/{user_id}")
async def get_user(user_id: str, user: dict = Depends(get_current_user)):
    if user["id"] != user_id:
        require_permission(user, "members", "view")
    u = await db.users.find_one({"id": user_id}, {"password_hash": 0, "_id": 0})
    if not u:
        raise HTTPException(404, "Not found")
    return u


@api.put("/users/{user_id}")
async def update_user(user_id: str, body: UserUpdate, request: Request, user: dict = Depends(get_current_user)):
    is_self = user["id"] == user_id
    if not is_self:
        require_permission(user, "members", "update")
    existing = await db.users.find_one({"id": user_id})
    if not existing:
        raise HTTPException(404, "Not found")
    # SECURITY: Developer account is immutable to non-developers
    if existing.get("role") == "developer" and user.get("role") != "developer":
        raise HTTPException(403, "Developer account cannot be modified by non-developer role")
    upd = {k: v for k, v in body.model_dump().items() if v is not None}
    # Only admins can change role/permissions/status; developers preserve their own protections
    if not has_permission(user, "members", "manage"):
        upd.pop("role", None)
        upd.pop("permissions", None)
        upd.pop("status", None)
    # Even with manage permission, only a developer can grant/revoke developer role
    if upd.get("role") == "developer" and user.get("role") != "developer":
        raise HTTPException(403, "Only a developer can assign developer role")
    upd["updated_at"] = now_iso()
    await db.users.update_one({"id": user_id}, {"$set": upd})
    await audit(user, "UPDATE", "members", user_id, {k: existing.get(k) for k in upd}, upd,
                request.client.host if request.client else "")
    u = await db.users.find_one({"id": user_id}, {"password_hash": 0, "_id": 0})
    return u


@api.delete("/users/{user_id}")
async def delete_user(user_id: str, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "members", "delete")
    if user_id == user["id"]:
        raise HTTPException(400, "Cannot delete self")
    target = await db.users.find_one({"id": user_id})
    if target and target.get("role") == "developer" and user.get("role") != "developer":
        raise HTTPException(403, "Developer account cannot be deleted by non-developer role")
    await db.users.update_one({"id": user_id}, {"$set": {"status": "disabled", "updated_at": now_iso()}})
    await audit(user, "DELETE", "members", user_id, None, None, request.client.host if request.client else "")
    return {"ok": True}


# ================= CUSTOMERS =================
@api.get("/customers")
async def list_customers(q: Optional[str] = None, status: Optional[str] = None,
                         user: dict = Depends(get_current_user)):
    require_permission(user, "customers", "view")
    query = {}
    if q:
        query["$or"] = [
            {"company_name": {"$regex": q, "$options": "i"}},
            {"contact_person": {"$regex": q, "$options": "i"}},
        ]
    if status:
        query["status"] = status
    docs = await db.customers.find(query, {"_id": 0}).to_list(500)
    # Auto-inactive if contract ended
    today = date.today().isoformat()
    for d in docs:
        if d.get("contract_end") and d["contract_end"] < today and d.get("status") != "inactive":
            await db.customers.update_one({"id": d["id"]}, {"$set": {"status": "inactive"}})
            d["status"] = "inactive"
    return docs


@api.post("/customers")
async def create_customer(body: CustomerCreate, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "customers", "create")
    doc = body.model_dump()
    # Remove client credentials from customer doc
    client_email = doc.pop("client_email", None)
    client_password = doc.pop("client_password", None)
    doc.update({
        "id": uid(),
        "status": "active",
        "registration_date": now_iso(),
        "last_visit": None,
        "created_by": user["id"],
        "created_at": now_iso(),
        "updated_at": now_iso(),
    })
    await db.customers.insert_one(doc)
    # Optionally create a client login account tied to this customer
    if client_email and client_password:
        client_email = client_email.lower()
        if not await db.users.find_one({"email": client_email}):
            await db.users.insert_one({
                "id": uid(),
                "email": client_email,
                "password_hash": hash_password(client_password),
                "full_name": doc.get("contact_person") or doc["company_name"],
                "role": "client",
                "position": "Client",
                "phone": doc.get("phone", ""),
                "address": doc.get("address", ""),
                "customer_id": doc["id"],
                "id_number": "",
                "leave_quota": 0,
                "leave_used": 0,
                "status": "active",
                "profile_photo": None,
                "permissions": default_client_permissions(),
                "created_at": now_iso(),
                "updated_at": now_iso(),
            })
    await audit(user, "CREATE", "customers", doc["id"], None, {"company_name": doc["company_name"]}, "")
    return clean(doc)


@api.put("/customers/{cid}")
async def update_customer(cid: str, body: CustomerCreate, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "customers", "update")
    existing = await db.customers.find_one({"id": cid})
    if not existing:
        raise HTTPException(404, "Not found")
    upd = body.model_dump()
    upd["updated_at"] = now_iso()
    await db.customers.update_one({"id": cid}, {"$set": upd})
    await audit(user, "UPDATE", "customers", cid, {"company_name": existing.get("company_name")}, upd, "")
    return await db.customers.find_one({"id": cid}, {"_id": 0})


@api.delete("/customers/{cid}")
async def delete_customer(cid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "customers", "delete")
    await db.customers.update_one({"id": cid}, {"$set": {"status": "inactive"}})
    await audit(user, "DELETE", "customers", cid)
    return {"ok": True}


# ================= TASKS =================
def compute_task_status(task: dict) -> str:
    if task.get("status") == "cancelled":
        return "cancelled"
    if task.get("service_report_id") and task.get("check_out_at"):
        return "completed"
    if task.get("check_in_at"):
        return "in_progress"
    # Overdue check
    sched = task.get("scheduled_date")
    dl = task.get("deadline")
    today = date.today().isoformat()
    if dl and dl < now_iso():
        return "overdue"
    if sched and sched < today:
        return "overdue"
    return "pending"


@api.get("/tasks")
async def list_tasks(filter: Optional[str] = None,
                     customer_id: Optional[str] = None,
                     technician_id: Optional[str] = None,
                     date_from: Optional[str] = None,
                     date_to: Optional[str] = None,
                     user: dict = Depends(get_current_user)):
    require_permission(user, "tasks", "view")
    query = {}
    # Client role: only see their assigned customer (hard-enforced, overrides any query param)
    if user.get("role") == "client":
        if customer_id and customer_id != user.get("customer_id"):
            raise HTTPException(403, "Forbidden: cannot query other customer's data")
        query["customer_id"] = user.get("customer_id") or "__none__"
    elif user.get("role") == "technician" and not has_permission(user, "tasks", "manage"):
        query["technician_id"] = user["id"]
        if customer_id:
            query["customer_id"] = customer_id
    else:
        if customer_id:
            query["customer_id"] = customer_id
    if technician_id and user.get("role") != "client":
        query["technician_id"] = technician_id
    if date_from:
        query["scheduled_date"] = {"$gte": date_from}
    if date_to:
        query.setdefault("scheduled_date", {})["$lte"] = date_to
    docs = await db.tasks.find(query, {"_id": 0}).sort("scheduled_date", -1).to_list(1000)
    # Enrich + compute status
    cust_ids = list({d["customer_id"] for d in docs})
    tech_ids = list({d["technician_id"] for d in docs})
    cmap = {c["id"]: c for c in await db.customers.find({"id": {"$in": cust_ids}}, {"_id": 0}).to_list(500)}
    tmap = {t["id"]: t for t in await db.users.find({"id": {"$in": tech_ids}}, {"_id": 0, "password_hash": 0}).to_list(500)}
    result = []
    for d in docs:
        d["status"] = compute_task_status(d)
        d["customer"] = cmap.get(d["customer_id"])
        d["technician"] = tmap.get(d["technician_id"])
        if filter == "pending" and d["status"] != "pending":
            continue
        if filter == "overdue" and d["status"] != "overdue":
            continue
        if filter == "completed" and d["status"] != "completed":
            continue
        if filter == "in_progress" and d["status"] != "in_progress":
            continue
        result.append(d)
    return result


@api.post("/tasks")
async def create_task(body: TaskCreate, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "tasks", "create")
    doc = body.model_dump()
    doc.update({
        "id": uid(),
        "assigned_by": user["id"],
        "status": "pending",
        "check_in_at": None,
        "check_out_at": None,
        "service_report_id": None,
        "created_at": now_iso(),
        "updated_at": now_iso(),
    })
    await db.tasks.insert_one(doc)
    await db.notifications.insert_one({
        "id": uid(),
        "user_id": body.technician_id,
        "title": "New Task Assigned",
        "message": f"You have been assigned a new task: {body.work_target}",
        "type": "task",
        "record_id": doc["id"],
        "read": False,
        "created_at": now_iso(),
    })
    await audit(user, "CREATE", "tasks", doc["id"], None, {"target": body.work_target}, "")
    return clean(doc)


@api.get("/tasks/{tid}")
async def get_task(tid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "tasks", "view")
    t = await db.tasks.find_one({"id": tid}, {"_id": 0})
    if not t:
        raise HTTPException(404, "Not found")
    if user.get("role") != "admin" and t["technician_id"] != user["id"] and not has_permission(user, "tasks", "manage"):
        raise HTTPException(403, "Forbidden")
    t["status"] = compute_task_status(t)
    t["customer"] = await db.customers.find_one({"id": t["customer_id"]}, {"_id": 0})
    t["technician"] = await db.users.find_one({"id": t["technician_id"]}, {"_id": 0, "password_hash": 0})
    return t


@api.put("/tasks/{tid}")
async def update_task(tid: str, body: dict = Body(...), user: dict = Depends(get_current_user)):
    require_permission(user, "tasks", "update")
    body.pop("id", None)
    body["updated_at"] = now_iso()
    await db.tasks.update_one({"id": tid}, {"$set": body})
    await audit(user, "UPDATE", "tasks", tid, None, body)
    return await db.tasks.find_one({"id": tid}, {"_id": 0})


@api.delete("/tasks/{tid}")
async def delete_task(tid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "tasks", "delete")
    await db.tasks.delete_one({"id": tid})
    await audit(user, "DELETE", "tasks", tid)
    return {"ok": True}


@api.post("/tasks/{tid}/reopen")
async def reopen_task(tid: str, user: dict = Depends(get_current_user)):
    """Admin or Developer can reopen a completed task — clears service_report_id
    and reverts status so the technician can re-submit."""
    if user.get("role") not in ("admin", "developer") and not has_permission(user, "tasks", "manage"):
        raise HTTPException(403, "Only admin/developer can reopen tasks")
    t = await db.tasks.find_one({"id": tid})
    if not t:
        raise HTTPException(404, "Not found")
    # Optionally soft-delete the old SR so audit trail keeps it
    if t.get("service_report_id"):
        await db.service_reports.update_one({"id": t["service_report_id"]}, {"$set": {"status": "reopened", "reopened_at": now_iso()}})
    await db.tasks.update_one({"id": tid}, {"$set": {
        "service_report_id": None,
        "check_out_at": None,
        "status": "pending",
        "updated_at": now_iso(),
    }})
    await audit(user, "REOPEN", "tasks", tid, {"prev_sr": t.get("service_report_id")}, None)
    return {"ok": True}


# ================= ATTENDANCE =================
@api.post("/attendance/checkin")
async def checkin(body: AttendanceCheckIn, user: dict = Depends(get_current_user)):
    require_permission(user, "attendance", "create")
    geofence_ok = True
    distance = None
    if body.task_id:
        task = await db.tasks.find_one({"id": body.task_id}, {"_id": 0})
        if task:
            cust = await db.customers.find_one({"id": task["customer_id"]})
            if cust and cust.get("latitude") and cust.get("longitude"):
                distance = haversine_m(body.latitude, body.longitude, cust["latitude"], cust["longitude"])
                radius = float(os.environ.get("GEOFENCE_RADIUS_METERS", "100"))
                geofence_ok = distance <= radius
    address = await reverse_geocode(body.latitude, body.longitude)
    doc = {
        "id": uid(),
        "user_id": user["id"],
        "user_name": user.get("full_name"),
        "task_id": body.task_id,
        "type": "check_in",
        "latitude": body.latitude,
        "longitude": body.longitude,
        "accuracy": body.accuracy,
        "photo": body.photo,
        "address": address,
        "geofence_ok": geofence_ok,
        "distance_meters": distance,
        "date": date.today().isoformat(),
        "timestamp": now_iso(),
    }
    await db.attendance.insert_one(doc)
    if body.task_id:
        await db.tasks.update_one({"id": body.task_id}, {"$set": {"check_in_at": doc["timestamp"], "updated_at": now_iso()}})
    return clean(doc)


@api.post("/attendance/checkout")
async def checkout(body: AttendanceCheckOut, user: dict = Depends(get_current_user)):
    require_permission(user, "attendance", "create")
    ci = await db.attendance.find_one({"id": body.attendance_id})
    if not ci:
        raise HTTPException(404, "Check-in not found")
    address = await reverse_geocode(body.latitude, body.longitude)
    # Compute working hours
    working_hours = None
    try:
        d1 = datetime.fromisoformat(ci["timestamp"])
        d2 = datetime.now(timezone.utc)
        working_hours = round((d2 - d1).total_seconds() / 3600.0, 2)
    except Exception:
        pass
    doc = {
        "id": uid(),
        "user_id": user["id"],
        "user_name": user.get("full_name"),
        "task_id": ci.get("task_id"),
        "type": "check_out",
        "latitude": body.latitude,
        "longitude": body.longitude,
        "accuracy": body.accuracy,
        "photo": body.photo,
        "address": address,
        "checkin_ref": body.attendance_id,
        "working_hours": working_hours,
        "date": date.today().isoformat(),
        "timestamp": now_iso(),
    }
    await db.attendance.insert_one(doc)
    # update the check-in row so listing shows working_hours + checkout_address on the pair
    await db.attendance.update_one({"id": body.attendance_id}, {"$set": {
        "checkout_id": doc["id"], "working_hours": working_hours,
        "checkout_address": address, "checkout_timestamp": doc["timestamp"],
    }})
    if ci.get("task_id"):
        await db.tasks.update_one({"id": ci["task_id"]}, {"$set": {"check_out_at": doc["timestamp"], "updated_at": now_iso()}})
    return clean(doc)


@api.get("/attendance")
async def list_attendance(user_id: Optional[str] = None,
                          customer_id: Optional[str] = None,
                          date_from: Optional[str] = None,
                          date_to: Optional[str] = None,
                          user: dict = Depends(get_current_user)):
    require_permission(user, "attendance", "view")
    query = {}
    if user.get("role") == "client":
        # Filter to attendance where task's customer = client's
        cid = user.get("customer_id") or "__none__"
        tasks = await db.tasks.find({"customer_id": cid}, {"_id": 0, "id": 1}).to_list(2000)
        query["task_id"] = {"$in": [t["id"] for t in tasks]}
    elif user_id:
        query["user_id"] = user_id
    elif not has_permission(user, "attendance", "manage") and user.get("role") not in ("admin",):
        query["user_id"] = user["id"]
    if customer_id and user.get("role") != "client":
        tasks = await db.tasks.find({"customer_id": customer_id}, {"_id": 0, "id": 1}).to_list(2000)
        query["task_id"] = {"$in": [t["id"] for t in tasks]}
    if date_from:
        query["date"] = {"$gte": date_from}
    if date_to:
        query.setdefault("date", {})["$lte"] = date_to
    docs = await db.attendance.find(query, {"_id": 0}).sort("timestamp", -1).to_list(1000)
    return docs


# ================= GPS TRACKING =================
@api.post("/gps/ping")
async def gps_ping(body: GPSPing, user: dict = Depends(get_current_user)):
    # Users track their own location — require travel.track OR own-tracking (admins always allowed)
    if user.get("role") != "admin" and not has_permission(user, "travel", "track"):
        raise HTTPException(status_code=403, detail="Missing permission: travel.track")
    doc = body.model_dump()
    doc.update({
        "id": uid(),
        "user_id": user["id"],
        "timestamp": now_iso(),
    })
    await db.gps_tracking.insert_one(doc)
    await db.users.update_one({"id": user["id"]}, {"$set": {
        "last_lat": body.latitude,
        "last_lng": body.longitude,
        "last_gps_at": doc["timestamp"],
        "last_seen": doc["timestamp"],
        "online": True,
        "gps_status": "active",
    }})
    return {"ok": True}


@api.post("/heartbeat")
async def heartbeat(body: HeartbeatBody, user: dict = Depends(get_current_user)):
    upd = {"last_seen": now_iso(), "online": True, "gps_status": body.gps_status}
    if body.latitude and body.longitude:
        upd["last_lat"] = body.latitude
        upd["last_lng"] = body.longitude
        upd["last_gps_at"] = now_iso()
    await db.users.update_one({"id": user["id"]}, {"$set": upd})
    return {"ok": True}


@api.get("/gps/tracks")
async def gps_tracks(user_id: Optional[str] = None, task_id: Optional[str] = None,
                     date_from: Optional[str] = None, date_to: Optional[str] = None,
                     user: dict = Depends(get_current_user)):
    require_permission(user, "location", "view")
    query = {}
    if user_id:
        query["user_id"] = user_id
    if task_id:
        query["task_id"] = task_id
    if date_from:
        query["timestamp"] = {"$gte": date_from}
    if date_to:
        query.setdefault("timestamp", {})["$lte"] = date_to + "T23:59:59"
    docs = await db.gps_tracking.find(query, {"_id": 0}).sort("timestamp", 1).limit(2000).to_list(2000)
    return docs


@api.get("/location/live")
async def live_locations(user: dict = Depends(get_current_user)):
    require_permission(user, "location", "view")
    query = {"last_lat": {"$exists": True}}
    if user.get("role") != "admin" and not has_permission(user, "location", "track"):
        query["id"] = user["id"]
    users = await db.users.find(query, {"_id": 0, "password_hash": 0}).to_list(200)
    # Compute online status: heartbeat within 90s
    threshold = datetime.now(timezone.utc) - timedelta(seconds=90)
    for u in users:
        try:
            ls = datetime.fromisoformat(u.get("last_seen", ""))
            u["online"] = ls >= threshold
        except Exception:
            u["online"] = False
    return users


# ================= TRAVEL / PERJALANAN =================
@api.get("/travel")
async def travel_summary(user_id: Optional[str] = None, date_from: Optional[str] = None,
                         date_to: Optional[str] = None, user: dict = Depends(get_current_user)):
    require_permission(user, "travel", "view")
    query = {}
    if user_id:
        query["user_id"] = user_id
    elif user.get("role") != "admin":
        query["user_id"] = user["id"]
    if date_from:
        query["timestamp"] = {"$gte": date_from}
    if date_to:
        query.setdefault("timestamp", {})["$lte"] = date_to + "T23:59:59"
    tracks = await db.gps_tracking.find(query, {"_id": 0}).sort("timestamp", 1).to_list(5000)

    # Group by user_id + task_id
    groups: Dict[str, List[dict]] = {}
    for t in tracks:
        key = f"{t['user_id']}::{t.get('task_id') or 'nomission'}"
        groups.setdefault(key, []).append(t)

    out = []
    for key, pts in groups.items():
        if len(pts) < 2:
            continue
        dist = 0.0
        for i in range(1, len(pts)):
            dist += haversine_m(pts[i - 1]["latitude"], pts[i - 1]["longitude"], pts[i]["latitude"], pts[i]["longitude"])
        u_id, t_id = key.split("::")
        out.append({
            "user_id": u_id,
            "task_id": None if t_id == "nomission" else t_id,
            "start_time": pts[0]["timestamp"],
            "end_time": pts[-1]["timestamp"],
            "distance_m": round(dist, 2),
            "point_count": len(pts),
            "start": {"lat": pts[0]["latitude"], "lng": pts[0]["longitude"]},
            "end": {"lat": pts[-1]["latitude"], "lng": pts[-1]["longitude"]},
        })
    return out


# ================= SERVICE REPORT =================
@api.post("/service-reports")
async def create_sr(body: ServiceReportCreate, request: Request, user: dict = Depends(get_current_user)):
    require_permission(user, "service_reports", "create")
    task = await db.tasks.find_one({"id": body.task_id})
    if not task:
        raise HTTPException(404, "Task not found")
    # Auto-generate report number
    count = await db.service_reports.count_documents({})
    report_no = f"SR-{datetime.now().strftime('%Y%m')}-{count + 1:04d}"
    now = datetime.now(timezone.utc)
    doc = body.model_dump()
    doc.update({
        "id": uid(),
        "report_number": report_no,
        "date": now.date().isoformat(),
        "time": now.strftime("%H:%M:%S"),
        "technician_id": task["technician_id"],
        "customer_id": task["customer_id"],
        "status": "submitted",
        "created_by": user["id"],
        "created_at": now_iso(),
        "updated_at": now_iso(),
    })
    await db.service_reports.insert_one(doc)
    await db.tasks.update_one({"id": body.task_id}, {"$set": {"service_report_id": doc["id"], "status": "completed", "updated_at": now_iso()}})
    await audit(user, "CREATE", "service_reports", doc["id"], None, {"task_id": body.task_id, "report": report_no})
    return clean(doc)


@api.get("/service-reports")
async def list_sr(customer_id: Optional[str] = None,
                  technician_id: Optional[str] = None,
                  date_from: Optional[str] = None,
                  date_to: Optional[str] = None,
                  user: dict = Depends(get_current_user)):
    require_permission(user, "service_reports", "view")
    query = {}
    if user.get("role") == "client":
        if customer_id and customer_id != user.get("customer_id"):
            raise HTTPException(403, "Forbidden")
        query["customer_id"] = user.get("customer_id") or "__none__"
    elif user.get("role") == "technician" and not has_permission(user, "service_reports", "manage"):
        query["technician_id"] = user["id"]
        if customer_id:
            query["customer_id"] = customer_id
    else:
        if customer_id:
            query["customer_id"] = customer_id
    if technician_id and user.get("role") != "client":
        query["technician_id"] = technician_id
    if date_from:
        query["date"] = {"$gte": date_from}
    if date_to:
        query.setdefault("date", {})["$lte"] = date_to
    docs = await db.service_reports.find(query, {"_id": 0}).sort("created_at", -1).to_list(1000)
    # Enrich with customer + technician names
    cust_ids = list({d["customer_id"] for d in docs})
    tech_ids = list({d["technician_id"] for d in docs})
    cmap = {c["id"]: c for c in await db.customers.find({"id": {"$in": cust_ids}}, {"_id": 0}).to_list(500)}
    tmap = {t["id"]: t for t in await db.users.find({"id": {"$in": tech_ids}}, {"_id": 0, "password_hash": 0}).to_list(500)}
    for d in docs:
        d["customer_name"] = cmap.get(d["customer_id"], {}).get("company_name", "")
        d["technician_name"] = tmap.get(d["technician_id"], {}).get("full_name", "")
    return docs


@api.get("/service-reports/{sid}")
async def get_sr(sid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "service_reports", "view")
    sr = await db.service_reports.find_one({"id": sid}, {"_id": 0})
    if not sr:
        raise HTTPException(404, "Not found")
    sr["task"] = await db.tasks.find_one({"id": sr["task_id"]}, {"_id": 0})
    sr["customer"] = await db.customers.find_one({"id": sr["customer_id"]}, {"_id": 0})
    sr["technician"] = await db.users.find_one({"id": sr["technician_id"]}, {"_id": 0, "password_hash": 0})
    return sr


@api.get("/service-reports/{sid}/pdf")
async def sr_pdf(sid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "service_reports", "export")
    sr = await db.service_reports.find_one({"id": sid}, {"_id": 0})
    if not sr:
        raise HTTPException(404, "Not found")
    task = await db.tasks.find_one({"id": sr["task_id"]}, {"_id": 0}) or {}
    cust = await db.customers.find_one({"id": sr["customer_id"]}, {"_id": 0}) or {}
    tech = await db.users.find_one({"id": sr["technician_id"]}, {"_id": 0, "password_hash": 0}) or {}
    # Client scope check
    if user.get("role") == "client" and user.get("customer_id") != sr["customer_id"]:
        raise HTTPException(403, "Forbidden")
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand = {
        "company_name": settings.get("company_name") or os.environ.get("COMPANY_NAME", "PestOps Pro"),
        "company_address": settings.get("company_address") or os.environ.get("COMPANY_ADDRESS", ""),
        "company_email": settings.get("company_email") or os.environ.get("COMPANY_EMAIL", ""),
    }
    if settings.get("logo_path"):
        try:
            data, _ = get_object(settings["logo_path"])
            brand["logo_bytes"] = data
        except Exception:
            pass
    # Fetch signatures + photos
    sig_bytes = {"photos": []}
    if sr.get("technician_signature"):
        try:
            sig_bytes["tech"], _ = get_object(sr["technician_signature"])
        except Exception:
            pass
    if sr.get("client_signature"):
        try:
            sig_bytes["client"], _ = get_object(sr["client_signature"])
        except Exception:
            pass
    for p in (sr.get("photos") or []):
        if isinstance(p, dict) and p.get("path"):
            try:
                pb, _ = get_object(p["path"])
                sig_bytes["photos"].append((pb, p.get("caption", "")))
            except Exception:
                pass
    pdf = generate_service_report_pdf(sr, task, cust, tech, brand=brand, sig_bytes=sig_bytes)
    return StreamingResponse(io.BytesIO(pdf), media_type="application/pdf",
                             headers={"Content-Disposition": f"attachment; filename={sr['report_number']}.pdf"})


# ================= LEAVE =================
@api.post("/leave")
async def create_leave(body: LeaveCreate, user: dict = Depends(get_current_user)):
    require_permission(user, "leave", "create")
    doc = body.model_dump()
    doc.update({
        "id": uid(),
        "user_id": user["id"],
        "user_name": user.get("full_name"),
        "status": "pending",
        "created_at": now_iso(),
        "reviewed_by": None,
        "reviewed_at": None,
    })
    await db.leave_requests.insert_one(doc)
    return clean(doc)


@api.get("/leave")
async def list_leave(user: dict = Depends(get_current_user)):
    require_permission(user, "leave", "view")
    query = {}
    if user.get("role") != "admin" and not has_permission(user, "leave", "approve"):
        query["user_id"] = user["id"]
    docs = await db.leave_requests.find(query, {"_id": 0}).sort("created_at", -1).to_list(500)
    return docs


@api.post("/leave/{lid}/decide")
async def decide_leave(lid: str, decision: str = Body(..., embed=True), user: dict = Depends(get_current_user)):
    require_permission(user, "leave", "approve")
    if decision not in ("approved", "rejected"):
        raise HTTPException(400, "Invalid decision")
    lv = await db.leave_requests.find_one({"id": lid})
    if not lv:
        raise HTTPException(404, "Not found")
    await db.leave_requests.update_one({"id": lid}, {"$set": {"status": decision, "reviewed_by": user["id"], "reviewed_at": now_iso()}})
    if decision == "approved":
        # Compute days
        try:
            sd = datetime.fromisoformat(lv["start_date"]).date()
            ed = datetime.fromisoformat(lv["end_date"]).date()
            days = (ed - sd).days + 1
            await db.users.update_one({"id": lv["user_id"]}, {"$inc": {"leave_used": days}})
        except Exception:
            pass
    await audit(user, decision.upper(), "leave", lid, {"status": "pending"}, {"status": decision})
    return {"ok": True}


# ================= FILE UPLOAD =================
@api.post("/upload")
async def upload_file(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    ext = (file.filename or "bin").rsplit(".", 1)[-1].lower()
    path = f"{APP_NAME}/uploads/{user['id']}/{uid()}.{ext}"
    data = await file.read()
    result = put_object(path, data, file.content_type or "application/octet-stream")
    await db.files.insert_one({
        "id": uid(),
        "storage_path": result["path"],
        "original_filename": file.filename,
        "content_type": file.content_type,
        "size": result.get("size", len(data)),
        "uploaded_by": user["id"],
        "created_at": now_iso(),
    })
    return {"path": result["path"], "size": result.get("size", len(data))}


@api.post("/upload/base64")
async def upload_b64(payload: dict = Body(...), user: dict = Depends(get_current_user)):
    """Upload base64 dataURL (used for camera photos & signatures)."""
    import base64
    data_url = payload.get("data", "")
    ext = payload.get("ext", "png")
    if "," in data_url:
        header, b64 = data_url.split(",", 1)
        # detect content type
        ct = "image/png"
        if "image/jpeg" in header:
            ct = "image/jpeg"
        elif "image/webp" in header:
            ct = "image/webp"
    else:
        b64 = data_url
        ct = "image/png"
    binary = base64.b64decode(b64)
    path = f"{APP_NAME}/uploads/{user['id']}/{uid()}.{ext}"
    result = put_object(path, binary, ct)
    return {"path": result["path"]}


@api.get("/files/{path:path}")
async def download_file(path: str, request: Request, auth: Optional[str] = Query(None)):
    # allow token via query param for <img> tags
    try:
        token = request.cookies.get("access_token") or auth
        if not token:
            raise HTTPException(401, "Not authenticated")
        from auth import decode_token
        decode_token(token)
    except Exception:
        raise HTTPException(401, "Invalid token")
    try:
        data, ct = get_object(path)
    except Exception:
        raise HTTPException(404, "File not found")
    return Response(content=data, media_type=ct)


# ================= NOTIFICATIONS =================
@api.get("/notifications")
async def list_notifs(user: dict = Depends(get_current_user)):
    docs = await db.notifications.find({"user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).limit(50).to_list(50)
    return docs


@api.post("/notifications/{nid}/read")
async def mark_read(nid: str, user: dict = Depends(get_current_user)):
    await db.notifications.update_one({"id": nid, "user_id": user["id"]}, {"$set": {"read": True}})
    return {"ok": True}


# ================= AUDIT LOG =================
@api.get("/audit-logs")
async def list_audit(user: dict = Depends(get_current_user)):
    require_permission(user, "audit_log", "view")
    docs = await db.audit_logs.find({}, {"_id": 0}).sort("timestamp", -1).limit(500).to_list(500)
    return docs


# ================= DASHBOARD =================
@api.get("/dashboard")
async def dashboard(user: dict = Depends(get_current_user)):
    is_admin = user.get("role") == "admin" or has_permission(user, "tasks", "manage")
    task_q = {} if is_admin else {"technician_id": user["id"]}
    all_tasks = await db.tasks.find(task_q, {"_id": 0}).to_list(1000)
    today = date.today().isoformat()
    for t in all_tasks:
        t["status"] = compute_task_status(t)
    tasks_summary = {
        "total": len(all_tasks),
        "pending": sum(1 for t in all_tasks if t["status"] == "pending"),
        "overdue": sum(1 for t in all_tasks if t["status"] == "overdue"),
        "completed": sum(1 for t in all_tasks if t["status"] == "completed"),
        "today": sum(1 for t in all_tasks if t.get("scheduled_date") == today),
    }

    # Technician summary (admin only)
    tech_summary = None
    if is_admin:
        techs = await db.users.find({"role": "technician"}, {"_id": 0, "password_hash": 0}).to_list(500)
        threshold = datetime.now(timezone.utc) - timedelta(seconds=90)
        online = 0
        for tt in techs:
            try:
                ls = datetime.fromisoformat(tt.get("last_seen", ""))
                if ls >= threshold:
                    online += 1
            except Exception:
                pass
        # On task = has task in progress
        in_prog_tech = set()
        for t in all_tasks:
            if t["status"] == "in_progress":
                in_prog_tech.add(t["technician_id"])
        tech_summary = {
            "total": len(techs),
            "online": online,
            "offline": len(techs) - online,
            "on_task": len(in_prog_tech),
            "not_on_task": len(techs) - len(in_prog_tech),
        }

    # Attendance today
    today_attend = await db.attendance.find({"date": today}, {"_id": 0}).to_list(500)
    checked_in = len({a["user_id"] for a in today_attend if a["type"] == "check_in"})
    checked_out = len({a["user_id"] for a in today_attend if a["type"] == "check_out"})
    attendance_summary = {"checked_in": checked_in, "checked_out": checked_out}

    # Pest findings this month
    month_prefix = datetime.now(timezone.utc).strftime("%Y-%m")
    srs = await db.service_reports.find({"date": {"$regex": f"^{month_prefix}"}}, {"_id": 0}).to_list(500)
    pest_totals = {"F": 0, "M": 0, "C": 0, "R": 0, "A": 0, "O": 0}
    for sr in srs:
        for f in sr.get("pest_findings") or []:
            if f.get("code") in pest_totals:
                pest_totals[f["code"]] += int(f.get("quantity") or 0)

    return {
        "tasks": tasks_summary,
        "technicians": tech_summary,
        "attendance": attendance_summary,
        "pest_findings_month": pest_totals,
    }


# ================= REPORTS =================
@api.get("/reports/attendance")
async def report_attendance(format: str = "excel", date_from: Optional[str] = None,
                            date_to: Optional[str] = None, month: Optional[str] = None,
                            user_id: Optional[str] = None, customer_id: Optional[str] = None,
                            user: dict = Depends(get_current_user)):
    require_permission(user, "reports", "export")
    q = {}
    if user_id:
        q["user_id"] = user_id
    if month:
        q["date"] = {"$regex": f"^{month}"}
    else:
        if date_from:
            q["date"] = {"$gte": date_from}
        if date_to:
            q.setdefault("date", {})["$lte"] = date_to
    if customer_id:
        tasks_list = await db.tasks.find({"customer_id": customer_id}, {"_id": 0, "id": 1}).to_list(5000)
        q["task_id"] = {"$in": [t["id"] for t in tasks_list]}
    att = await db.attendance.find(q, {"_id": 0}).sort("timestamp", 1).to_list(5000)
    cust_map = {c["id"]: c for c in await db.customers.find({}, {"_id": 0}).to_list(5000)}
    task_map = {t["id"]: t for t in await db.tasks.find({}, {"_id": 0}).to_list(5000)}
    rows = []
    for a in att:
        if a.get("type") != "check_in":
            continue
        co = next((x for x in att if x.get("checkin_ref") == a["id"]), None)
        wh = a.get("working_hours") or ""
        if co and not wh:
            try:
                d1 = datetime.fromisoformat(a["timestamp"]); d2 = datetime.fromisoformat(co["timestamp"])
                wh = f"{(d2 - d1).total_seconds() / 3600:.2f}"
            except Exception:
                pass
        cust_name = ""
        if a.get("task_id") and a["task_id"] in task_map:
            cust_name = cust_map.get(task_map[a["task_id"]]["customer_id"], {}).get("company_name", "")
        rows.append([
            a.get("user_name", ""), cust_name,
            a["date"], a["timestamp"][11:19],
            a.get("address", "") or f"{a.get('latitude', '')},{a.get('longitude', '')}",
            (co["date"] if co else ""), (co["timestamp"][11:19] if co else ""),
            (co.get("address", "") if co else ""),
            str(wh),
        ])
    headers = ["Employee", "Client", "Check-in Date", "Check-in Time", "Check-in Location",
               "Check-out Date", "Check-out Time", "Check-out Location", "Working Hours"]
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand = {"company_name": settings.get("company_name"), "company_address": settings.get("company_address"), "company_email": settings.get("company_email")}
    if format == "pdf":
        b = generate_simple_pdf("ATTENDANCE REPORT", headers, rows, brand=brand)
        return StreamingResponse(io.BytesIO(b), media_type="application/pdf",
                                 headers={"Content-Disposition": "attachment; filename=attendance.pdf"})
    b = generate_simple_excel("Attendance", headers, rows)
    return StreamingResponse(io.BytesIO(b),
                             media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=attendance.xlsx"})


@api.get("/reports/customers")
async def report_customers(format: str = "excel", date_from: Optional[str] = None,
                           date_to: Optional[str] = None, customer_id: Optional[str] = None,
                           user: dict = Depends(get_current_user)):
    require_permission(user, "reports", "export")
    q = {}
    if customer_id:
        q["id"] = customer_id
    if date_from:
        q["registration_date"] = {"$gte": date_from}
    if date_to:
        q.setdefault("registration_date", {})["$lte"] = date_to + "T23:59:59"
    docs = await db.customers.find(q, {"_id": 0}).to_list(5000)
    headers = ["Company", "Project", "Contact", "Phone", "Email", "Address", "Category", "Status", "Contract Start", "Contract End", "Registered"]
    rows = [[d.get("company_name", ""), d.get("project_name", ""), d.get("contact_person", ""), d.get("phone", ""),
             d.get("email", ""), d.get("address", ""), d.get("category", ""), d.get("status", ""),
             d.get("contract_start") or "", d.get("contract_end") or "", (d.get("registration_date") or "")[:10]] for d in docs]
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand = {"company_name": settings.get("company_name"), "company_address": settings.get("company_address"), "company_email": settings.get("company_email")}
    if format == "pdf":
        b = generate_simple_pdf("CUSTOMER REPORT", headers, rows, brand=brand)
        return StreamingResponse(io.BytesIO(b), media_type="application/pdf",
                                 headers={"Content-Disposition": "attachment; filename=customers.pdf"})
    b = generate_simple_excel("Customers", headers, rows)
    return StreamingResponse(io.BytesIO(b),
                             media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=customers.xlsx"})


@api.get("/reports/employees")
async def report_employees(format: str = "excel", date_from: Optional[str] = None,
                           date_to: Optional[str] = None, user_id: Optional[str] = None,
                           role: Optional[str] = None, user: dict = Depends(get_current_user)):
    require_permission(user, "reports", "export")
    q = {}
    if user_id:
        q["id"] = user_id
    if role:
        q["role"] = role
    if date_from:
        q["created_at"] = {"$gte": date_from}
    if date_to:
        q.setdefault("created_at", {})["$lte"] = date_to + "T23:59:59"
    docs = await db.users.find(q, {"_id": 0, "password_hash": 0}).to_list(5000)
    headers = ["Name", "Email", "Role", "Position", "Phone", "ID Number", "Address", "Join Date", "Status", "Leave Quota", "Remaining"]
    rows = [[d.get("full_name", ""), d.get("email", ""), d.get("role", ""), d.get("position", ""),
             d.get("phone", ""), d.get("id_number", ""), d.get("address", ""),
             (d.get("created_at", "") or "")[:10], d.get("status", ""),
             d.get("leave_quota", 0), max(0, (d.get("leave_quota", 0) or 0) - (d.get("leave_used", 0) or 0))] for d in docs]
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand = {"company_name": settings.get("company_name"), "company_address": settings.get("company_address"), "company_email": settings.get("company_email")}
    if format == "pdf":
        b = generate_simple_pdf("EMPLOYEE DATA REPORT", headers, rows, brand=brand)
        return StreamingResponse(io.BytesIO(b), media_type="application/pdf",
                                 headers={"Content-Disposition": "attachment; filename=employees.pdf"})
    b = generate_simple_excel("Employees", headers, rows)
    return StreamingResponse(io.BytesIO(b),
                             media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=employees.xlsx"})


# ================= SETTINGS =================
@api.get("/settings")
async def get_settings(user: dict = Depends(get_current_user)):
    doc = await db.settings.find_one({"_id": "app"}, {"_id": 0}) or {}
    doc.setdefault("company_name", os.environ.get("COMPANY_NAME", "Proteksi Pest Control"))
    doc.setdefault("company_address", os.environ.get("COMPANY_ADDRESS", ""))
    doc.setdefault("company_email", os.environ.get("COMPANY_EMAIL", ""))
    doc.setdefault("geofence_radius", int(os.environ.get("GEOFENCE_RADIUS_METERS", "100")))
    doc.setdefault("gps_interval", int(os.environ.get("GPS_INTERVAL_SECONDS", "4")))
    return doc


@api.put("/settings")
async def update_settings(body: dict = Body(...), user: dict = Depends(get_current_user)):
    require_permission(user, "settings", "manage")
    body["updated_at"] = now_iso()
    await db.settings.update_one({"_id": "app"}, {"$set": body}, upsert=True)
    return await get_settings(user)


# ================= Health =================
@api.get("/health")
async def health():
    try:
        await db.command("ping")
        return {"status": "ok", "db": "up", "timestamp": now_iso()}
    except Exception as e:
        return {"status": "degraded", "error": str(e)}



# ================= SCHEDULES (recurring standby) =================
@api.get("/schedules")
async def list_schedules(customer_id: Optional[str] = None,
                         technician_id: Optional[str] = None,
                         date_from: Optional[str] = None,
                         date_to: Optional[str] = None,
                         user: dict = Depends(get_current_user)):
    require_permission(user, "schedule", "view")
    query = {}
    if user.get("role") == "client":
        if customer_id and customer_id != user.get("customer_id"):
            raise HTTPException(403, "Forbidden")
        query["customer_id"] = user.get("customer_id") or "__none__"
    elif user.get("role") == "technician" and not has_permission(user, "schedule", "manage"):
        query["technician_id"] = user["id"]
        if customer_id:
            query["customer_id"] = customer_id
    else:
        if customer_id:
            query["customer_id"] = customer_id
    if technician_id and user.get("role") != "client":
        query["technician_id"] = technician_id
    if date_from:
        query["date"] = {"$gte": date_from}
    if date_to:
        query.setdefault("date", {})["$lte"] = date_to
    docs = await db.schedules.find(query, {"_id": 0}).sort("date", 1).to_list(2000)
    cust_ids = list({d["customer_id"] for d in docs})
    tech_ids = list({d["technician_id"] for d in docs})
    cmap = {c["id"]: c for c in await db.customers.find({"id": {"$in": cust_ids}}, {"_id": 0}).to_list(500)}
    tmap = {t["id"]: t for t in await db.users.find({"id": {"$in": tech_ids}}, {"_id": 0, "password_hash": 0}).to_list(500)}
    for d in docs:
        d["customer"] = cmap.get(d["customer_id"])
        d["technician"] = tmap.get(d["technician_id"])
    return docs


@api.post("/schedules/mass-create")
async def mass_create_schedules(body: ScheduleMassCreate, user: dict = Depends(get_current_user)):
    require_permission(user, "schedule", "create")
    try:
        sd = datetime.fromisoformat(body.start_date).date()
        ed = datetime.fromisoformat(body.end_date).date()
    except Exception:
        raise HTTPException(400, "Invalid dates")
    if ed < sd:
        raise HTTPException(400, "End date before start date")
    if (ed - sd).days > 366:
        raise HTTPException(400, "Range too large (max 366 days)")
    weekdays = set(body.weekdays or [0, 1, 2, 3, 4, 5])
    created = []
    cur = sd
    while cur <= ed:
        if cur.weekday() in weekdays:
            doc = {
                "id": uid(),
                "customer_id": body.customer_id,
                "technician_id": body.technician_id,
                "date": cur.isoformat(),
                "start_time": body.start_time,
                "end_time": body.end_time,
                "status": "scheduled",
                "notes": body.notes,
                "created_by": user["id"],
                "created_at": now_iso(),
                "updated_at": now_iso(),
            }
            created.append(doc)
        cur = cur + timedelta(days=1)
    if created:
        await db.schedules.insert_many(created)
    await audit(user, "CREATE", "schedule", f"mass-{len(created)}", None, {"count": len(created), "customer": body.customer_id, "technician": body.technician_id})
    return {"count": len(created), "schedules": [clean(d) for d in created]}


@api.put("/schedules/{sid}")
async def update_schedule(sid: str, body: ScheduleUpdate, user: dict = Depends(get_current_user)):
    require_permission(user, "schedule", "update")
    upd = {k: v for k, v in body.model_dump().items() if v is not None}
    upd["updated_at"] = now_iso()
    await db.schedules.update_one({"id": sid}, {"$set": upd})
    await audit(user, "UPDATE", "schedule", sid, None, upd)
    return await db.schedules.find_one({"id": sid}, {"_id": 0})


@api.delete("/schedules/{sid}")
async def delete_schedule(sid: str, user: dict = Depends(get_current_user)):
    require_permission(user, "schedule", "delete")
    await db.schedules.delete_one({"id": sid})
    await audit(user, "DELETE", "schedule", sid)
    return {"ok": True}


# ================= GEOCODING (Nominatim proxy) =================
@api.get("/geocode/search")
async def geocode_search(q: str, user: dict = Depends(get_current_user)):
    if not q or len(q) < 3:
        return []
    try:
        async with httpx.AsyncClient(timeout=10) as cli:
            r = await cli.get("https://nominatim.openstreetmap.org/search",
                              params={"q": q, "format": "json", "limit": 5, "addressdetails": 1},
                              headers={"User-Agent": "PestOpsPro/1.0"})
            r.raise_for_status()
            data = r.json()
        return [{"display_name": d.get("display_name"),
                 "lat": float(d.get("lat", 0)), "lng": float(d.get("lon", 0))} for d in data]
    except Exception as e:
        logger.warning(f"Geocode failed: {e}")
        return []


@api.get("/geocode/reverse")
async def geocode_reverse_api(lat: float, lon: float, user: dict = Depends(get_current_user)):
    return {"display_name": await reverse_geocode(lat, lon)}


# ================= MONTHLY REPORT =================
@api.get("/monthly-report")
async def monthly_report(customer_id: str, month: str, user: dict = Depends(get_current_user)):
    """month = YYYY-MM. Returns aggregated data + PDF-ready payload."""
    require_permission(user, "monthly_reports", "view")
    if user.get("role") == "client" and user.get("customer_id") != customer_id:
        raise HTTPException(403, "Forbidden")
    customer = await db.customers.find_one({"id": customer_id}, {"_id": 0})
    if not customer:
        raise HTTPException(404, "Customer not found")
    year, mon = month.split("-")
    m_prefix = f"{year}-{mon}"
    target = date(int(year), int(mon), 1)
    # Historical anchor: FIRST service report ever for this customer (fallback to contract_start / target month)
    first_sr = await db.service_reports.find_one({"customer_id": customer_id}, sort=[("date", 1)])
    if first_sr and first_sr.get("date"):
        try:
            cs = datetime.fromisoformat(first_sr["date"]).date().replace(day=1)
        except Exception:
            cs = date(int(year), 1, 1)
    else:
        contract_start = customer.get("contract_start") or f"{year}-01-01"
        try:
            cs = datetime.fromisoformat(contract_start).date().replace(day=1)
        except Exception:
            cs = date(int(year), 1, 1)
    # Ensure target month is always included even if there are no SRs yet
    if cs > target:
        cs = target

    # Historical pest per-month from contract_start → target
    historical = []
    cur = cs
    while cur <= target:
        prefix = cur.strftime("%Y-%m")
        srs_m = await db.service_reports.find(
            {"customer_id": customer_id, "date": {"$regex": f"^{prefix}"}}, {"_id": 0}
        ).to_list(500)
        totals = {"F": 0, "M": 0, "C": 0, "R": 0, "A": 0, "O": 0}
        for sr in srs_m:
            for f in sr.get("pest_findings") or []:
                if f.get("code") in totals:
                    totals[f["code"]] += int(f.get("quantity") or 0)
        historical.append({"month": prefix, **totals, "total": sum(totals.values())})
        # next month
        y2, m2 = (cur.year + (1 if cur.month == 12 else 0), 1 if cur.month == 12 else cur.month + 1)
        cur = date(y2, m2, 1)

    # Current month reports/attendance/work
    srs = await db.service_reports.find({"customer_id": customer_id, "date": {"$regex": f"^{m_prefix}"}}, {"_id": 0}).sort("date", 1).to_list(500)
    tasks = await db.tasks.find({"customer_id": customer_id, "scheduled_date": {"$regex": f"^{m_prefix}"}}, {"_id": 0}).to_list(500)
    task_ids = [t["id"] for t in tasks]
    attendance = await db.attendance.find({"task_id": {"$in": task_ids}, "type": "check_in"}, {"_id": 0}).sort("timestamp", 1).to_list(500)

    # Enrich SRs with tech names
    tech_ids = list({s["technician_id"] for s in srs})
    tmap = {t["id"]: t for t in await db.users.find({"id": {"$in": tech_ids}}, {"_id": 0, "password_hash": 0}).to_list(500)}
    for s in srs:
        s["technician_name"] = tmap.get(s["technician_id"], {}).get("full_name", "")

    return {
        "customer": customer,
        "month": m_prefix,
        "contract_start": customer.get("contract_start") or cs.isoformat(),
        "first_report_date": (first_sr.get("date") if first_sr else None),
        "historical_pest": historical,
        "service_reports": srs,
        "tasks": tasks,
        "attendance": attendance,
    }


def _generate_monthly_pdf(payload: dict, brand: dict) -> bytes:
    from reportlab.platypus import Image
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=15 * mm, bottomMargin=15 * mm)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph(f"<b>{brand.get('company_name','PestOps')}</b>", styles["Title"]))
    story.append(Paragraph(brand.get('company_address', ''), styles["Normal"]))
    story.append(Paragraph(f"<b>MONTHLY REPORT — {payload['month']}</b>", styles["Heading2"]))
    story.append(Spacer(1, 5 * mm))

    c = payload["customer"]
    story.append(Paragraph("<b>CLIENT INFORMATION</b>", styles["Heading4"]))
    ct = Table([["Company", c.get("company_name", "")], ["Project", c.get("project_name", "") or c.get("company_name", "")],
                ["Address", c.get("address", "")], ["Contract Start", c.get("contract_start", "")]],
               colWidths=[35 * mm, 140 * mm])
    ct.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey), ("FONTSIZE", (0, 0), (-1, -1), 9)]))
    story.append(ct)
    story.append(Spacer(1, 4 * mm))

    # Attendance table
    story.append(Paragraph("<b>EMPLOYEE ATTENDANCE</b>", styles["Heading4"]))
    if payload["attendance"]:
        rows = [["Technician", "Date", "Check-in", "Check-out", "Working Hours"]]
        for a in payload["attendance"]:
            rows.append([a.get("user_name", ""), a.get("date", ""), a.get("timestamp", "")[11:19],
                         (a.get("checkout_timestamp", "") or "")[11:19] if a.get("checkout_timestamp") else "-",
                         str(a.get("working_hours") or "-")])
        t = Table(rows, repeatRows=1)
        t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                               ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2937")),
                               ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                               ("FONTSIZE", (0, 0), (-1, -1), 8)]))
        story.append(t)
    else:
        story.append(Paragraph("No attendance records this period.", styles["Normal"]))
    story.append(Spacer(1, 4 * mm))

    # Work summary from SRs
    story.append(Paragraph("<b>WORK REALIZATION</b>", styles["Heading4"]))
    if payload["service_reports"]:
        rows = [["Date", "Technician", "Scope", "Recommendation"]]
        for s in payload["service_reports"]:
            rows.append([s.get("date", ""), s.get("technician_name", ""),
                         (s.get("scope_of_area", "") or "")[:40],
                         (s.get("recommendation", "") or "")[:60]])
        t = Table(rows, repeatRows=1, colWidths=[22 * mm, 35 * mm, 55 * mm, 63 * mm])
        t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                               ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2937")),
                               ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                               ("FONTSIZE", (0, 0), (-1, -1), 8)]))
        story.append(t)
    story.append(Spacer(1, 4 * mm))

    # Pest historical chart (as table since chart is on frontend)
    story.append(Paragraph("<b>PEST FINDINGS — HISTORICAL (contract start → current month)</b>", styles["Heading4"]))
    rows = [["Month", "F", "M", "C", "R", "A", "O", "Total"]]
    for h in payload["historical_pest"]:
        rows.append([h["month"], str(h["F"]), str(h["M"]), str(h["C"]), str(h["R"]), str(h["A"]), str(h["O"]), str(h["total"])])
    t = Table(rows, repeatRows=1)
    t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                           ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F2937")),
                           ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                           ("FONTSIZE", (0, 0), (-1, -1), 8)]))
    story.append(t)
    story.append(Spacer(1, 4 * mm))

    # Photo documentation (thumbnails names only for compact PDF)
    all_photos = []
    for s in payload["service_reports"]:
        for p in (s.get("photos") or []):
            if isinstance(p, dict) and p.get("path"):
                all_photos.append(p)
    if all_photos:
        story.append(Paragraph("<b>PHOTO DOCUMENTATION</b>", styles["Heading4"]))
        rows = [["#", "Caption", "Path"]]
        for i, p in enumerate(all_photos, 1):
            rows.append([str(i), (p.get("caption") or "")[:50], p["path"][-40:]])
        t = Table(rows, repeatRows=1, colWidths=[10 * mm, 80 * mm, 85 * mm])
        t.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                               ("FONTSIZE", (0, 0), (-1, -1), 8)]))
        story.append(t)

    doc.build(story)
    return buf.getvalue()


@api.get("/monthly-report/pdf")
async def monthly_report_pdf(customer_id: str, month: str, include_srs: bool = True,
                             user: dict = Depends(get_current_user)):
    require_permission(user, "monthly_reports", "export")
    payload = await monthly_report(customer_id=customer_id, month=month, user=user)
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand = {
        "company_name": settings.get("company_name") or os.environ.get("COMPANY_NAME", "PestOps Pro"),
        "company_address": settings.get("company_address") or os.environ.get("COMPANY_ADDRESS", ""),
        "company_email": settings.get("company_email") or os.environ.get("COMPANY_EMAIL", ""),
    }
    if settings.get("logo_path"):
        try:
            data, _ = get_object(settings["logo_path"])
            brand["logo_bytes"] = data
        except Exception:
            pass
    mp = _generate_monthly_pdf(payload, brand)
    if not include_srs or not payload["service_reports"]:
        return StreamingResponse(io.BytesIO(mp), media_type="application/pdf",
                                 headers={"Content-Disposition": f"attachment; filename=monthly-{customer_id[:8]}-{month}.pdf"})
    # Merge monthly + each SR PDF using pypdf
    try:
        from pypdf import PdfWriter, PdfReader
    except ImportError:
        logger.warning("pypdf not installed — Monthly Report SR merge disabled, returning summary-only PDF")
        return StreamingResponse(io.BytesIO(mp), media_type="application/pdf",
                                 headers={"Content-Disposition": f"attachment; filename=monthly-{customer_id[:8]}-{month}.pdf"})
    writer = PdfWriter()
    writer.append(PdfReader(io.BytesIO(mp)))
    for sr in payload["service_reports"]:
        task = await db.tasks.find_one({"id": sr["task_id"]}, {"_id": 0}) or {}
        cust = payload["customer"]
        tech = await db.users.find_one({"id": sr["technician_id"]}, {"_id": 0, "password_hash": 0}) or {}
        sig_bytes = {"photos": []}
        if sr.get("technician_signature"):
            try: sig_bytes["tech"], _ = get_object(sr["technician_signature"])
            except Exception: pass
        if sr.get("client_signature"):
            try: sig_bytes["client"], _ = get_object(sr["client_signature"])
            except Exception: pass
        for p in (sr.get("photos") or []):
            if isinstance(p, dict) and p.get("path"):
                try:
                    pb, _ = get_object(p["path"])
                    sig_bytes["photos"].append((pb, p.get("caption", "")))
                except Exception:
                    pass
        srpdf = generate_service_report_pdf(sr, task, cust, tech, brand=brand, sig_bytes=sig_bytes)
        writer.append(PdfReader(io.BytesIO(srpdf)))
    out = io.BytesIO()
    writer.write(out)
    return StreamingResponse(io.BytesIO(out.getvalue()), media_type="application/pdf",
                             headers={"Content-Disposition": f"attachment; filename=monthly-{customer_id[:8]}-{month}-full.pdf"})


@api.get("/monthly-report/excel")
async def monthly_report_excel(customer_id: str, month: str, user: dict = Depends(get_current_user)):
    require_permission(user, "monthly_reports", "export")
    payload = await monthly_report(customer_id=customer_id, month=month, user=user)
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    c = payload["customer"]
    ws.append(["MONTHLY REPORT", payload["month"]])
    ws.append(["Company", c.get("company_name", "")])
    ws.append(["Project", c.get("project_name", "")])
    ws.append(["Address", c.get("address", "")])
    ws.append([])
    ws.append(["Historical Pest Findings"])
    ws.append(["Month", "Fly (F)", "Mosquito (M)", "Cockroach (C)", "Rodent (R)", "Ant (A)", "Other (O)", "Total"])
    for h in payload["historical_pest"]:
        ws.append([h["month"], h["F"], h["M"], h["C"], h["R"], h["A"], h["O"], h["total"]])
    ws2 = wb.create_sheet("Work Realization")
    ws2.append(["Date", "Technician", "Scope", "Recommendation", "Pest Description"])
    for s in payload["service_reports"]:
        ws2.append([s.get("date", ""), s.get("technician_name", ""),
                    s.get("scope_of_area", ""), s.get("recommendation", ""),
                    s.get("pest_description", "")])
    ws3 = wb.create_sheet("Attendance")
    ws3.append(["Technician", "Date", "Check-in", "Check-out", "Working Hours", "Address"])
    for a in payload["attendance"]:
        ws3.append([a.get("user_name", ""), a.get("date", ""),
                    (a.get("timestamp") or "")[11:19],
                    (a.get("checkout_timestamp") or "")[11:19] or "-",
                    a.get("working_hours") or "-", a.get("address", "")])
    buf = io.BytesIO(); wb.save(buf)
    return StreamingResponse(io.BytesIO(buf.getvalue()),
                             media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": f"attachment; filename=monthly-{customer_id[:8]}-{month}.xlsx"})


@api.get("/monthly-report/pptx")
async def monthly_report_pptx(customer_id: str, month: str, user: dict = Depends(get_current_user)):
    require_permission(user, "monthly_reports", "export")
    payload = await monthly_report(customer_id=customer_id, month=month, user=user)
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    brand = await _get_brand()
    prs = Presentation()
    # Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "MONTHLY REPORT"
    s1.placeholders[1].text = f"{payload['customer'].get('company_name','')} — {payload['month']}"
    # Client info
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Client Information"
    c = payload["customer"]
    tf = s2.placeholders[1].text_frame
    tf.text = f"Company: {c.get('company_name','')}"
    for line in [f"Project: {c.get('project_name','') or '-'}",
                 f"Address: {c.get('address','')}",
                 f"Contract Start: {c.get('contract_start','') or '-'}",
                 f"Period: {payload['month']}"]:
        p = tf.add_paragraph(); p.text = line
    # Historical
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Pest Findings — Historical"
    left = Inches(0.5); top = Inches(1.4); width = Inches(9); height = Inches(4)
    rows = len(payload["historical_pest"]) + 1
    tbl = s3.shapes.add_table(rows, 8, left, top, width, height).table
    for j, h in enumerate(["Month", "F", "M", "C", "R", "A", "O", "Total"]):
        tbl.cell(0, j).text = h
    for i, h in enumerate(payload["historical_pest"], start=1):
        tbl.cell(i, 0).text = h["month"]
        for j, k in enumerate(["F", "M", "C", "R", "A", "O", "total"], start=1):
            tbl.cell(i, j).text = str(h[k])
    # Work Realization
    s4 = prs.slides.add_slide(prs.slide_layouts[5])
    s4.shapes.title.text = f"Work Realization ({payload['month']})"
    if payload["service_reports"]:
        rows = min(len(payload["service_reports"]) + 1, 12)
        tbl2 = s4.shapes.add_table(rows, 3, Inches(0.5), Inches(1.4), Inches(9), Inches(4)).table
        for j, h in enumerate(["Date", "Technician", "Scope"]):
            tbl2.cell(0, j).text = h
        for i, s in enumerate(payload["service_reports"][:rows - 1], start=1):
            tbl2.cell(i, 0).text = s.get("date", "")
            tbl2.cell(i, 1).text = s.get("technician_name", "")
            tbl2.cell(i, 2).text = (s.get("scope_of_area", "") or "")[:60]
    # Attendance
    s5 = prs.slides.add_slide(prs.slide_layouts[5])
    s5.shapes.title.text = "Attendance"
    if payload["attendance"]:
        rows = min(len(payload["attendance"]) + 1, 15)
        tbl3 = s5.shapes.add_table(rows, 4, Inches(0.5), Inches(1.4), Inches(9), Inches(4.5)).table
        for j, h in enumerate(["Technician", "Date", "Check-in", "Hours"]):
            tbl3.cell(0, j).text = h
        for i, a in enumerate(payload["attendance"][:rows - 1], start=1):
            tbl3.cell(i, 0).text = a.get("user_name", "")
            tbl3.cell(i, 1).text = a.get("date", "")
            tbl3.cell(i, 2).text = (a.get("timestamp") or "")[11:19]
            tbl3.cell(i, 3).text = str(a.get("working_hours") or "-")
    buf = io.BytesIO(); prs.save(buf)
    return StreamingResponse(io.BytesIO(buf.getvalue()),
                             media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                             headers={"Content-Disposition": f"attachment; filename=monthly-{customer_id[:8]}-{month}.pptx"})


# ================= BULK ZIP EXPORT =================
@api.get("/service-reports/export/zip")
async def sr_bulk_zip(customer_id: Optional[str] = None,
                      date_from: Optional[str] = None,
                      date_to: Optional[str] = None,
                      user: dict = Depends(get_current_user)):
    require_permission(user, "service_reports", "export")
    query = {}
    if user.get("role") == "client":
        query["customer_id"] = user.get("customer_id") or "__none__"
    elif customer_id:
        query["customer_id"] = customer_id
    if date_from:
        query["date"] = {"$gte": date_from}
    if date_to:
        query.setdefault("date", {})["$lte"] = date_to
    srs = await db.service_reports.find(query, {"_id": 0}).to_list(500)
    if not srs:
        raise HTTPException(404, "No reports found")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for sr in srs:
            task = await db.tasks.find_one({"id": sr["task_id"]}, {"_id": 0}) or {}
            cust = await db.customers.find_one({"id": sr["customer_id"]}, {"_id": 0}) or {}
            tech = await db.users.find_one({"id": sr["technician_id"]}, {"_id": 0, "password_hash": 0}) or {}
            pdf = generate_service_report_pdf(sr, task, cust, tech)
            zf.writestr(f"{sr.get('report_number', sr['id'])}.pdf", pdf)
    buf.seek(0)
    return StreamingResponse(io.BytesIO(buf.getvalue()), media_type="application/zip",
                             headers={"Content-Disposition": "attachment; filename=service_reports.zip"})


# ================= EMAIL SEND =================
async def _get_brand():
    s = await db.settings.find_one({"_id": "app"}) or {}
    return {
        "company_name": s.get("company_name") or os.environ.get("COMPANY_NAME", "PestOps Pro"),
        "company_address": s.get("company_address") or os.environ.get("COMPANY_ADDRESS", ""),
    }


async def _email_config() -> dict:
    return await db.settings.find_one({"_id": "email"}) or {}


def _mask_pw(cfg: dict) -> dict:
    """Remove SMTP password & Twilio token from responses; expose only whether they are set."""
    out = {k: v for k, v in cfg.items() if k not in ("_id", "smtp_password", "wa_auth_token")}
    out["smtp_password_set"] = bool(cfg.get("smtp_password"))
    out["wa_auth_token_set"] = bool(cfg.get("wa_auth_token"))
    return out


async def _smart_send(*, to: str, subject: str, plain_body: str, attachments: list = None):
    """Try Custom SMTP first (if configured), else fallback to Emergent Resend.
    Raises HTTPException if email disabled globally."""
    cfg = await _email_config()
    if cfg.get("email_enabled") is False:
        raise HTTPException(status_code=409, detail="Email delivery disabled in settings")
    brand = await _get_brand()
    html = body_to_html(plain_body, signature=cfg.get("signature", ""))

    smtp_host = cfg.get("smtp_host")
    smtp_user = cfg.get("smtp_username")
    smtp_pw = cfg.get("smtp_password")
    from_addr = cfg.get("from_email")

    if smtp_host and smtp_user and smtp_pw and from_addr:
        try:
            return await send_via_smtp(
                host=smtp_host, port=int(cfg.get("smtp_port") or 587),
                username=smtp_user, password=smtp_pw,
                use_tls=bool(cfg.get("smtp_use_tls", True)),
                from_addr=from_addr, from_name=cfg.get("from_name") or brand["company_name"],
                to=to, subject=subject, html=html,
                attachments=attachments or [],
                reply_to=cfg.get("reply_to"),
            )
        except Exception as e:
            logger.error(f"SMTP send failed, falling back to Resend: {e}")

    # Fallback to Emergent Resend
    return await send_email_with_attachments(
        to=to, subject=subject, html=html,
        attachments=attachments or [], reply_to=cfg.get("reply_to"),
    )


@api.get("/email-settings")
async def get_email_settings(user: dict = Depends(get_current_user)):
    if user.get("role") not in ("admin", "developer") and not has_permission(user, "settings", "manage"):
        raise HTTPException(403, "Forbidden")
    cfg = await _email_config()
    cfg.setdefault("smtp_host", "")
    cfg.setdefault("smtp_port", 587)
    cfg.setdefault("smtp_username", "")
    cfg.setdefault("smtp_use_tls", True)
    cfg.setdefault("from_email", "")
    cfg.setdefault("from_name", "")
    cfg.setdefault("reply_to", "")
    cfg.setdefault("signature", "")
    cfg.setdefault("sr_subject_template", DEFAULT_SR_SUBJECT)
    cfg.setdefault("sr_body_template", DEFAULT_SR_BODY)
    cfg.setdefault("mr_subject_template", DEFAULT_MR_SUBJECT)
    cfg.setdefault("mr_body_template", DEFAULT_MR_BODY)
    cfg.setdefault("auto_monthly_send", False)
    cfg.setdefault("auto_monthly_day", 1)  # send on 1st of month
    cfg.setdefault("email_enabled", True)
    # WhatsApp defaults
    cfg.setdefault("wa_enabled", False)
    cfg.setdefault("wa_account_sid", "")
    cfg.setdefault("wa_from", os.environ.get("TWILIO_WA_FROM", "whatsapp:+14155238886"))
    cfg.setdefault("wa_sr_template", DEFAULT_SR_WA)
    cfg.setdefault("wa_mr_template", DEFAULT_MR_WA)
    cfg.setdefault("wa_auto_monthly", False)
    return _mask_pw(cfg)


@api.put("/email-settings")
async def put_email_settings(body: dict = Body(...), user: dict = Depends(get_current_user)):
    if user.get("role") not in ("admin", "developer") and not has_permission(user, "settings", "manage"):
        raise HTTPException(403, "Forbidden")
    # Whitelist keys (no arbitrary fields)
    ALLOWED = {"smtp_host", "smtp_port", "smtp_username", "smtp_password", "smtp_use_tls",
               "from_email", "from_name", "reply_to", "signature",
               "sr_subject_template", "sr_body_template",
               "mr_subject_template", "mr_body_template",
               "auto_monthly_send", "auto_monthly_day", "email_enabled",
               # WhatsApp
               "wa_enabled", "wa_account_sid", "wa_auth_token", "wa_from",
               "wa_sr_template", "wa_mr_template", "wa_auto_monthly"}
    upd = {k: v for k, v in body.items() if k in ALLOWED}
    # Don't overwrite password/token with empty or null (keep existing)
    for secret_key in ("smtp_password", "wa_auth_token"):
        if secret_key in upd and (upd[secret_key] in ("", None)):
            upd.pop(secret_key)
    upd["updated_at"] = now_iso()
    await db.settings.update_one({"_id": "email"}, {"$set": upd}, upsert=True)
    await audit(user, "UPDATE", "settings", "email", None, {k: v for k, v in upd.items() if k != "smtp_password"})
    return await get_email_settings(user)


@api.post("/email-settings/test")
async def test_email(body: dict = Body(...), user: dict = Depends(get_current_user)):
    if user.get("role") not in ("admin", "developer"):
        raise HTTPException(403, "Forbidden")
    to = body.get("to") or user.get("email")
    if not to:
        raise HTTPException(400, "No recipient")
    brand = await _get_brand()
    plain = f"This is a test email from {brand['company_name']} — PestOps Pro platform.\n\nIf you received this, your email integration is working correctly.\n\nSent by: {user.get('full_name')}"
    try:
        sent = await _smart_send(to=to, subject=f"Test Email — {brand['company_name']}",
                                 plain_body=plain, attachments=[])
        return {"ok": True, "sent_via": sent}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Email send failed: {str(e)}")


# ================= WHATSAPP =================
async def _wa_send(to: str, body: str) -> str:
    cfg = await _email_config()
    if not cfg.get("wa_enabled"):
        raise HTTPException(status_code=409, detail="WhatsApp disabled in settings")
    sid = cfg.get("wa_account_sid") or os.environ.get("TWILIO_ACCOUNT_SID", "")
    token = cfg.get("wa_auth_token") or os.environ.get("TWILIO_AUTH_TOKEN", "")
    wa_from = cfg.get("wa_from") or os.environ.get("TWILIO_WA_FROM", "")
    if not (sid and token and wa_from):
        raise HTTPException(status_code=400, detail="WhatsApp credentials missing")
    return await send_whatsapp(account_sid=sid, auth_token=token, from_wa=wa_from,
                               to=to, body=body)


@api.post("/wa/test")
async def test_wa(body: dict = Body(...), user: dict = Depends(get_current_user)):
    if user.get("role") not in ("admin", "developer"):
        raise HTTPException(403, "Forbidden")
    to = body.get("to")
    if not to:
        raise HTTPException(400, "recipient 'to' required (e.g. +6281234567890)")
    brand = await _get_brand()
    msg = body.get("message") or f"Test WhatsApp from {brand['company_name']} — PestOps Pro."
    try:
        sid = await _wa_send(to, msg)
        return {"ok": True, "sid": sid}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"WhatsApp send failed: {e}")


class SendWAReport(BaseModel):
    to: Optional[str] = None
    message: Optional[str] = None


@api.post("/service-reports/{sid}/whatsapp")
async def wa_single_sr(sid: str, body: SendWAReport, user: dict = Depends(get_current_user)):
    require_permission(user, "email", "create")
    sr = await db.service_reports.find_one({"id": sid}, {"_id": 0})
    if not sr:
        raise HTTPException(404, "Not found")
    cust = await db.customers.find_one({"id": sr["customer_id"]}, {"_id": 0}) or {}
    to = body.to or cust.get("phone")
    if not to:
        raise HTTPException(400, "No WhatsApp number on customer; provide 'to'.")
    cfg = await _email_config()
    brand = await _get_brand()
    ctx = {"client_name": cust.get("company_name", ""), "period": sr.get("date", ""),
           "report_number": sr.get("report_number", ""), "company_name": brand["company_name"],
           "technician": ""}
    msg = body.message or render_template(cfg.get("wa_sr_template", DEFAULT_SR_WA), **ctx)
    try:
        twilio_sid = await _wa_send(to, msg)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"WhatsApp send failed: {e}")
    await audit(user, "WHATSAPP", "service_reports", sid, None, {"to": to, "twilio_sid": twilio_sid})
    return {"ok": True, "sid": twilio_sid, "recipient": to}


@api.post("/monthly-report/whatsapp")
async def wa_monthly(customer_id: str = Body(...), month: str = Body(...),
                     body: SendWAReport = Body(...), user: dict = Depends(get_current_user)):
    require_permission(user, "email", "create")
    cust = await db.customers.find_one({"id": customer_id}, {"_id": 0})
    if not cust:
        raise HTTPException(404, "Customer not found")
    to = body.to or cust.get("phone")
    if not to:
        raise HTTPException(400, "No WhatsApp number available")
    cfg = await _email_config()
    brand = await _get_brand()
    ctx = {"client_name": cust.get("company_name", ""), "period": month, "company_name": brand["company_name"]}
    msg = body.message or render_template(cfg.get("wa_mr_template", DEFAULT_MR_WA), **ctx)
    try:
        twilio_sid = await _wa_send(to, msg)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"WhatsApp send failed: {e}")
    await audit(user, "WHATSAPP", "monthly_reports", f"{customer_id}-{month}", None, {"to": to, "twilio_sid": twilio_sid})
    return {"ok": True, "sid": twilio_sid, "recipient": to}


# ================= CSV IMPORT =================
@api.post("/customers/import-csv")
async def customers_import_csv(file: UploadFile = File(...), user: dict = Depends(get_current_user)):
    require_permission(user, "customers", "create")
    if not (file.filename or "").lower().endswith(".csv"):
        raise HTTPException(400, "CSV file required")
    raw = (await file.read()).decode("utf-8-sig", errors="ignore")
    reader = csv.DictReader(_io.StringIO(raw))
    created, skipped, errors = 0, 0, []
    for i, row in enumerate(reader, start=2):
        try:
            name = (row.get("company_name") or row.get("Company") or "").strip()
            if not name:
                skipped += 1
                continue
            if await db.customers.find_one({"company_name": name}):
                skipped += 1
                continue
            lat, lng = None, None
            try:
                lat = float(row["latitude"]) if row.get("latitude") else None
                lng = float(row["longitude"]) if row.get("longitude") else None
            except ValueError:
                pass
            doc = {
                "id": uid(),
                "company_name": name,
                "project_name": row.get("project_name", "") or row.get("project", ""),
                "contact_person": row.get("contact_person", "") or row.get("contact", ""),
                "phone": row.get("phone", ""),
                "email": row.get("email", ""),
                "address": row.get("address", ""),
                "location_text": row.get("location_text", ""),
                "latitude": lat, "longitude": lng,
                "category": row.get("category", "Regular"),
                "contract_start": row.get("contract_start") or None,
                "contract_end": row.get("contract_end") or None,
                "status": "active",
                "registration_date": now_iso(),
                "created_by": user["id"],
                "created_at": now_iso(), "updated_at": now_iso(),
            }
            await db.customers.insert_one(doc)
            created += 1
        except Exception as e:
            errors.append({"row": i, "error": str(e)[:120]})
    await audit(user, "IMPORT", "customers", f"csv-{created}", None, {"created": created, "skipped": skipped, "errors": len(errors)})
    return {"created": created, "skipped": skipped, "errors": errors}


@api.get("/customers/import-template.csv")
async def customers_import_template(user: dict = Depends(get_current_user)):
    csv_txt = "company_name,project_name,contact_person,phone,email,address,latitude,longitude,category,contract_start,contract_end\n"
    csv_txt += "PT. Contoh,Head Office,Budi,+628123456789,contoh@example.com,Jl. Sudirman Jakarta,-6.2088,106.8456,Corporate,2026-01-01,2026-12-31\n"
    return Response(content=csv_txt, media_type="text/csv",
                    headers={"Content-Disposition": "attachment; filename=customers_template.csv"})


@api.post("/service-reports/{sid}/email")
async def email_single_sr(sid: str, body: SendReportEmail, user: dict = Depends(get_current_user)):
    require_permission(user, "email", "create")
    sr = await db.service_reports.find_one({"id": sid}, {"_id": 0})
    if not sr:
        raise HTTPException(404, "Not found")
    cust = await db.customers.find_one({"id": sr["customer_id"]}, {"_id": 0}) or {}
    recipient = body.override_recipient or cust.get("email")
    if not recipient:
        raise HTTPException(400, "No client email available")
    task = await db.tasks.find_one({"id": sr["task_id"]}, {"_id": 0}) or {}
    tech = await db.users.find_one({"id": sr["technician_id"]}, {"_id": 0, "password_hash": 0}) or {}
    brand = await _get_brand()
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand_full = {**brand, "company_email": settings.get("company_email")}
    if settings.get("logo_path"):
        try:
            data, _ = get_object(settings["logo_path"])
            brand_full["logo_bytes"] = data
        except Exception:
            pass
    # Signatures + photos
    sig_bytes = {"photos": []}
    if sr.get("technician_signature"):
        try: sig_bytes["tech"], _ = get_object(sr["technician_signature"])
        except Exception: pass
    if sr.get("client_signature"):
        try: sig_bytes["client"], _ = get_object(sr["client_signature"])
        except Exception: pass
    for p in (sr.get("photos") or []):
        if isinstance(p, dict) and p.get("path"):
            try:
                pb, _ = get_object(p["path"])
                sig_bytes["photos"].append((pb, p.get("caption", "")))
            except Exception:
                pass
    pdf = generate_service_report_pdf(sr, task, cust, tech, brand=brand_full, sig_bytes=sig_bytes)
    # Templates
    cfg = await _email_config()
    ctx = {
        "report_number": sr.get("report_number", ""),
        "client_name": cust.get("company_name", ""),
        "period": sr.get("date", ""),
        "company_name": brand["company_name"],
        "technician": tech.get("full_name", ""),
    }
    subject = body.subject or render_template(cfg.get("sr_subject_template", DEFAULT_SR_SUBJECT), **ctx)
    plain = body.message or render_template(cfg.get("sr_body_template", DEFAULT_SR_BODY), **ctx)
    sent = await _smart_send(to=recipient, subject=subject, plain_body=plain,
                             attachments=[{"filename": f"{sr.get('report_number', 'report')}.pdf", "content": pdf}])
    await audit(user, "EMAIL", "service_reports", sid, None, {"to": recipient, "sent_via": sent})
    return {"ok": True, "sent_via": sent, "recipient": recipient}


@api.post("/monthly-report/email")
async def email_monthly_report(customer_id: str = Body(...), month: str = Body(...),
                               body: SendReportEmail = Body(...), user: dict = Depends(get_current_user)):
    require_permission(user, "email", "create")
    payload = await monthly_report(customer_id=customer_id, month=month, user=user)
    cust = payload["customer"]
    recipient = body.override_recipient or cust.get("email")
    if not recipient:
        raise HTTPException(400, "No client email available")
    brand = await _get_brand()
    settings = await db.settings.find_one({"_id": "app"}) or {}
    brand_full = {**brand, "company_email": settings.get("company_email")}
    if settings.get("logo_path"):
        try:
            data, _ = get_object(settings["logo_path"])
            brand_full["logo_bytes"] = data
        except Exception:
            pass
    pdf = _generate_monthly_pdf(payload, brand_full)
    cfg = await _email_config()
    ctx = {"client_name": cust.get("company_name", ""), "period": month, "company_name": brand["company_name"]}
    subject = body.subject or render_template(cfg.get("mr_subject_template", DEFAULT_MR_SUBJECT), **ctx)
    plain = body.message or render_template(cfg.get("mr_body_template", DEFAULT_MR_BODY), **ctx)
    sent = await _smart_send(to=recipient, subject=subject, plain_body=plain,
                             attachments=[{"filename": f"monthly-{month}.pdf", "content": pdf}])
    await audit(user, "EMAIL", "monthly_reports", f"{customer_id}-{month}", None, {"to": recipient, "sent_via": sent})
    return {"ok": True, "sent_via": sent, "recipient": recipient}


@api.post("/cron/auto-monthly-send")
async def cron_auto_monthly(request: Request):
    # Cron endpoints must ack 2xx immediately; enqueue/background the actual work.
    auth = request.headers.get("Authorization", "")
    expected = f"Bearer {os.environ.get('WEBHOOK_CRON_SECRET', '')}"
    if not expected.endswith(" ") and auth != expected:
        raise HTTPException(401, "Unauthorized")
    cfg = await _email_config()
    do_email = cfg.get("auto_monthly_send") and cfg.get("email_enabled", True)
    do_wa = cfg.get("wa_auto_monthly") and cfg.get("wa_enabled")
    if not do_email and not do_wa:
        return {"ok": True, "skipped": "email/wa auto disabled"}
    # Compute previous month
    today = date.today()
    if today.month == 1:
        prev_year, prev_month = today.year - 1, 12
    else:
        prev_year, prev_month = today.year, today.month - 1
    period = f"{prev_year:04d}-{prev_month:02d}"
    # Iterate active customers with email set — enqueue async, respond fast
    import asyncio

    async def _job():
        customers = await db.customers.find({"status": "active"}, {"_id": 0}).to_list(500)
        for c in customers:
            try:
                # Reuse email endpoint logic by simulating "admin"
                admin = await db.users.find_one({"role": "admin"}, {"_id": 0, "password_hash": 0})
                if not admin:
                    continue
                payload = await monthly_report(customer_id=c["id"], month=period, user=admin)
                brand = await _get_brand()
                settings = await db.settings.find_one({"_id": "app"}) or {}
                brand_full = {**brand, "company_email": settings.get("company_email")}
                if settings.get("logo_path"):
                    try:
                        data, _ = get_object(settings["logo_path"])
                        brand_full["logo_bytes"] = data
                    except Exception:
                        pass
                pdf = _generate_monthly_pdf(payload, brand_full)
                ctx = {"client_name": c.get("company_name", ""), "period": period, "company_name": brand["company_name"]}
                if do_email and c.get("email"):
                    subject = render_template(cfg.get("mr_subject_template", DEFAULT_MR_SUBJECT), **ctx)
                    plain = render_template(cfg.get("mr_body_template", DEFAULT_MR_BODY), **ctx)
                    try:
                        await _smart_send(to=c["email"], subject=subject, plain_body=plain,
                                          attachments=[{"filename": f"monthly-{period}.pdf", "content": pdf}])
                    except Exception as ee:
                        logger.warning(f"Auto-email {c.get('company_name')} failed: {ee}")
                if do_wa and c.get("phone"):
                    try:
                        msg = render_template(cfg.get("wa_mr_template", DEFAULT_MR_WA), **ctx)
                        await _wa_send(c["phone"], msg)
                    except Exception as we:
                        logger.warning(f"Auto-WA {c.get('company_name')} failed: {we}")
                await db.audit_logs.insert_one({
                    "id": uid(), "user_id": "cron", "user_name": "Scheduler",
                    "action": "AUTO_EMAIL", "module": "monthly_reports",
                    "record_id": f"{c['id']}-{period}", "new_value": {"to": c["email"]},
                    "timestamp": now_iso(),
                })
            except Exception as e:
                logger.exception(f"auto-monthly-send for {c.get('company_name')}: {e}")

    asyncio.create_task(_job())
    return {"ok": True, "period": period, "status": "queued"}



@api.get("/branding")
async def get_branding(user: dict = Depends(get_current_user)):
    s = await db.settings.find_one({"_id": "app"}) or {}
    return {
        "company_name": s.get("company_name") or "PestOps Pro",
        "company_address": s.get("company_address") or "",
        "company_email": s.get("company_email") or "",
        "company_phone": s.get("company_phone") or "",
        "logo_path": s.get("logo_path") or None,
        "app_name": s.get("app_name") or "PestOps Pro",
    }


@api.put("/branding")
async def put_branding(body: dict = Body(...), user: dict = Depends(get_current_user)):
    # SECURITY: Branding is Developer-only unless explicit permission granted
    if user.get("role") != "developer" and not has_permission(user, "branding", "manage"):
        raise HTTPException(403, "Branding is developer-only")
    body["updated_at"] = now_iso()
    await db.settings.update_one({"_id": "app"}, {"$set": body}, upsert=True)
    await audit(user, "UPDATE", "branding", "app", None, body)
    return await get_branding(user)



app.include_router(api)

# CORS - permissive for preview environment, cookies allowed
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,  # using Bearer token from localStorage as primary
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("shutdown")
async def _sd():
    client.close()
